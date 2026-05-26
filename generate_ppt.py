from pptx import Presentation
from pptx.util import Pt

def create_presentation():
    prs = Presentation()

    # Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Powered Itemized Kitchen Labels"
    subtitle.text = "Leveraging Star SMCS + StarIO.Online\n\nPresenter: [Your Name]\nStar Micronics – Board Presentation"

    # --- Slide 2: The Industry Problem ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Industry Problem"
    tf = slide.placeholders[1].text_frame
    tf.text = "Restaurants require itemized kitchen labels. But today:"
    tf.add_paragraph().text = "Square prints one format"
    tf.add_paragraph().text = "Uber Eats prints another"
    tf.add_paragraph().text = "DoorDash prints another"
    tf.add_paragraph().text = "Grubhub prints another"
    tf.add_paragraph().text = "No standardized structured order feed"
    tf.add_paragraph().text = "Result: Kitchen staff manually interpret thermal receipts."
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker Notes:\nRestaurants are increasingly requesting itemized labels for prep accuracy, allergy handling, and station routing. However, integration complexity across platforms blocks adoption."

    # --- Slide 3: Why This Is Difficult ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Why This Is Difficult"
    tf = slide.placeholders[1].text_frame
    tf.text = "Each POS platform:"
    tf.add_paragraph().text = "Uses different formatting & print commands"
    tf.add_paragraph().text = "Does not expose standardized structured order APIs"
    tf.add_paragraph().text = "Often restricts integration access"
    p = tf.add_paragraph()
    p.text = "\nTraditional integration approach:"
    tf.add_paragraph().text = "❌ Requires direct API integration with each POS"
    tf.add_paragraph().text = "❌ High maintenance cost & Long certification cycles"
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker Notes:\nThis is why most hardware manufacturers haven’t solved this problem yet."

    # --- Slide 4: The Insight ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Insight"
    tf = slide.placeholders[1].text_frame
    tf.text = "Star already owns the receipt layer."
    tf.add_paragraph().text = "Instead of integrating at the POS API layer… we standardize at the printed receipt image layer."
    p = tf.add_paragraph()
    p.text = "\n• Every platform prints."
    tf.add_paragraph().text = "• Every receipt can be uploaded."
    tf.add_paragraph().text = "• Every receipt can be interpreted."

    # --- Slide 5: System Architecture ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "System Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "POS / Delivery Platform"
    tf.add_paragraph().text = "↓  Star Receipt Printer"
    tf.add_paragraph().text = "↓  Receipt Image → SMCS"
    tf.add_paragraph().text = "↓  AI Receipt Parsing Engine"
    tf.add_paragraph().text = "↓  Structured Order JSON"
    tf.add_paragraph().text = "↓  StarIO.Online"
    tf.add_paragraph().text = "↓  Star Kitchen Label Printer"
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker Notes:\nWe are not competing with POS platforms. We are enhancing every POS platform."

    # --- Slide 6: Proof of Concept Demonstration ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Proof of Concept Demonstration"
    tf = slide.placeholders[1].text_frame
    tf.text = "The working system:"
    tf.add_paragraph().text = "Authenticates via SMCS OAuth"
    tf.add_paragraph().text = "Detects new receipt uploads"
    tf.add_paragraph().text = "Parses receipt image using AI"
    tf.add_paragraph().text = "Generates itemized label markup"
    tf.add_paragraph().text = "Prints to Star cloud-connected kitchen printer"
    tf.add_paragraph().text = "Handles crash recovery & duplicate prevention"

    # --- Slide 7: Why This Matters to Star ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Why This Matters to Star"
    tf = slide.placeholders[1].text_frame
    tf.text = "Strategic Advantages:"
    tf.add_paragraph().text = "✅ Platform-Agnostic (No POS API dependency)"
    tf.add_paragraph().text = "✅ Drives kitchen printer sales"
    tf.add_paragraph().text = "✅ Creates recurring SaaS model"
    tf.add_paragraph().text = "✅ Differentiates from Epson & Bixolon"
    tf.add_paragraph().text = "✅ Positions Star as AI-forward"
    p = tf.add_paragraph()
    p.text = "\nThis transforms Star from hardware vendor → intelligent infrastructure provider."

    # --- Slide 8: Monetization Opportunities ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Monetization Opportunities"
    tf = slide.placeholders[1].text_frame
    tf.text = "Potential business models:"
    tf.add_paragraph().text = "Monthly SaaS per location"
    tf.add_paragraph().text = "Per-label processing fee"
    tf.add_paragraph().text = "Premium SMCS subscription tier"
    tf.add_paragraph().text = "Enterprise kitchen automation package"
    tf.add_paragraph().text = "Bundled hardware + cloud subscription"

    # --- Slide 9: Competitive Advantage ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Competitive Advantage"
    tf = slide.placeholders[1].text_frame
    tf.text = "No competitor currently offers:"
    tf.add_paragraph().text = "Cloud-based AI receipt parsing"
    tf.add_paragraph().text = "Platform-agnostic kitchen label generation"
    tf.add_paragraph().text = "Integrated hardware + AI solution"
    p = tf.add_paragraph()
    p.text = "\nStar already owns:"
    tf.add_paragraph().text = "Hardware, Cloud infrastructure, Device authentication, Print API"

    # --- Slide 10: Risk & Mitigation ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Risk & Mitigation"
    tf = slide.placeholders[1].text_frame
    tf.text = "Risks:"
    tf.add_paragraph().text = "AI parsing variability & Receipt format edge cases"
    tf.add_paragraph().text = "Internet dependency"
    p = tf.add_paragraph()
    p.text = "\nMitigations:"
    tf.add_paragraph().text = "Continuous model tuning"
    tf.add_paragraph().text = "Fallback manual print"
    tf.add_paragraph().text = "Local caching & retry logic"
    tf.add_paragraph().text = "Edge inference in future hardware"

    # --- Slide 11: Expansion Potential ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Expansion Potential"
    tf = slide.placeholders[1].text_frame
    tf.text = "Future enhancements:"
    tf.add_paragraph().text = "Allergy auto-highlighting"
    tf.add_paragraph().text = "Rush order detection"
    tf.add_paragraph().text = "Prep time analytics"
    tf.add_paragraph().text = "Station routing"
    tf.add_paragraph().text = "Multi-language labels & QR codes"
    tf.add_paragraph().text = "Order-level analytics dashboard"

    # --- Slide 12: Strategic Positioning & Final Message ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Strategic Positioning"
    tf = slide.placeholders[1].text_frame
    tf.text = "Star Micronics can become:"
    tf.add_paragraph().text = "✅ The AI layer for hospitality printing"
    tf.add_paragraph().text = "✅ The infrastructure layer for POS-agnostic automation"
    tf.add_paragraph().text = "✅ A subscription-driven hardware ecosystem"
    p = tf.add_paragraph()
    p.text = "\nStar already owns the physical layer. This adds intelligence. This is a strategic evolution of Star’s cloud ecosystem."

    # --- Slide 13: Live Demonstration ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Live Demonstration"
    tf = slide.placeholders[1].text_frame
    tf.text = "Flow:"
    tf.add_paragraph().text = "1. Receipt prints"
    tf.add_paragraph().text = "2. Receipt uploads to SMCS"
    tf.add_paragraph().text = "3. AI parses receipt"
    tf.add_paragraph().text = "4. Labels print automatically"
    tf.add_paragraph().text = "5. Dashboard updates"

    # --- Slide 14: Appendix - Q&A ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Appendix: Anticipated Q&A"
    tf = slide.placeholders[1].text_frame
    tf.text = "Q: How accurate is the AI parsing?"
    tf.add_paragraph().text = "A: Accuracy is high in POC and improves over time. Fallback safeguards exist."
    tf.add_paragraph().text = "Q: What prevents competitors from doing this?"
    tf.add_paragraph().text = "A: Star already controls the hardware, cloud, and print API. Competitors lack this ecosystem."
    tf.add_paragraph().text = "Q: What about data privacy?"
    tf.add_paragraph().text = "A: Receipts are already in SMCS. Processing is controlled. No new exposure footprint."
    tf.add_paragraph().text = "Q: What is the revenue potential?"
    tf.add_paragraph().text = "A: Modest SaaS pricing across existing installed base yields meaningful recurring revenue."

    # Save the presentation
    filename = "Star_Micronics_AI_Kitchen_Labels.pptx"
    prs.save(filename)
    print(f"✅ Presentation successfully generated: {filename}")

if __name__ == "__main__":
    create_presentation()